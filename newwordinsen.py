
import docx

doc = docx.Document('Comp word list.docx')

#z1= doc.paragraphs[0].text
z2=[]
cnt=0
while  cnt < len(doc.paragraphs) :
    z2.append(doc.paragraphs[cnt].text)
    cnt=cnt+1
    
    
def noblank(x):
    return x!=''
def noblank1(x):
    return x!=' '

y2=list(filter(noblank,z2))

pp=list(range(0,len(y2)))

pp1=[]
for u in pp:
    pp1.append("WORD LIST " + str(u))
    
pp2=pp1[1:500]  

for n in pp2:
    if n in y2:
        y2.remove(n)
        
    else:
        continue
y3=[]
for i in y2:
    lelo=list(i)
    rere=list(filter(noblank1,lelo))  
    rere1="".join(rere)
    y3.append(rere1)
         

import re
y4=[]
for i in range(0,3172):
     review=re.sub('[^a-zA-Z]' , '' , y3[i])
     y4.append(review)








from urllib.error import HTTPError





from bs4 import BeautifulSoup as soup

from urllib.request import urlopen as uReq
from urllib.request import Request

doc2 = docx.Document()
#z2.remove('ambience')
count12=0
for l in y4:
        try:
            myurl=('https://wordsinasentence.com/' + l +'-in-a-sentence/' )
            
            req = Request(myurl, headers={'User-Agent': 'Mozilla/5.0'})
            
            uClient = uReq(req)
        
        except HTTPError as err:
            continue
        
        page_html = uClient.read()
        
        uClient.close()
        
        page_soup = soup(page_html,"html.parser")
        
        #containers = page_soup.findAll("div" ,{"class":"thecontent"})
        containers = page_soup.findAll("p")
        
        print(len(containers))
        a1= soup.prettify(containers[1])
        
        container=containers[0]
        #a2=container.findAll("div" , {"class":"thecontent clearfix"})
        #a3=a2[0].text
        a4= container.text
        
        eg1=[]
        for i in  containers:
            c1=i.text
            eg1.append(c1)
            
        
        
        doc2.add_paragraph(eg1[0])
        doc2.add_paragraph(eg1[1])
        ko=0
        for t in eg1:
            if (ko<2):
                print(ko)
            else:
                doc2.add_paragraph(eg1[ko])
            ko= ko+1  
            print(ko)
            
        print('hi this is this number of loop' +str(count12))
        
        
        
        doc2.save('word_in_sen_fin_12.docx')
        count12=count12+1
w1=[] 
count121=0      
for l in y4:
        try:
            myurl=('https://wordsinasentence.com/' + l +'-in-a-sentence/' )
            
            req = Request(myurl, headers={'User-Agent': 'Mozilla/5.0'})
            
            uClient = uReq(req)
        
        except HTTPError as err:
            continue
        
        
        page_html = uClient.read()
        
        uClient.close()
        
        page_soup = soup(page_html,"html.parser")
        
        #containers = page_soup.findAll("div" ,{"class":"thecontent"})
        containers = page_soup.findAll("p")
        
        print(len(containers))
        a1= soup.prettify(containers[1])
        
        container=containers[0]
        #a2=container.findAll("div" , {"class":"thecontent clearfix"})
        #a3=a2[0].text
        a4= container.text
        
        eg1=[]
        for i in  containers:
            c1=i.text
            eg1.append(c1)
        w2=[]
        w2.append(l) 
        hoho= eg1[1]+ "\n" + eg1[3]
        w2.append(hoho)
        w2.append(0)
        w1.append(w2)
        print('hi this is this number of loop' +str(count121))
        count121=count121+1
            
            






























        
from pptx import Presentation   

prs= Presentation() 
class Myslide:
    def __init__(self,data):
        self.title_slide_layout =prs.slide_layouts[data[2]]
        
        self.slide = prs.slides.add_slide(self.title_slide_layout)
        
        self.title = self.slide.shapes.title
        
        self.subtitle = self.slide.placeholders[1]
        
        self.title.text=data[0]
        
        self.subtitle.text=data[1] 


#slides=[["i ama dede","no\n  nono ",0] , ["i ama dele","no yo \n nono ",0]]    

for t in w1:
    
   Myslide(t)    
           

prs.save('pptword_fin51.pptx')
        
        


        
        
        
        
        

